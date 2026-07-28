import io
from pathlib import Path

import pytest
from docx import Document as DocxDocument
from fastapi.testclient import TestClient
from pptx import Presentation
from pypdf import PdfWriter

import app.api.documents as documents_module
from app.core.config import settings
from app.main import app
from app.services.embedding_service import EmbeddingService
from app.services.vector_store_service import VectorStoreService
from tests.fakes import fake_encoder


@pytest.fixture(autouse=True)
def isolated_upload_directory(tmp_path: Path, monkeypatch: pytest.MonkeyPatch) -> Path:
    """Redirect uploads to a throwaway directory so tests never touch the real uploads/ folder."""
    monkeypatch.setattr(settings, "upload_directory", tmp_path)
    return tmp_path


@pytest.fixture(autouse=True)
def isolated_vector_services(
    tmp_path: Path, monkeypatch: pytest.MonkeyPatch
) -> VectorStoreService:
    """
    Replace the real embedding and vector-store services with fast test
    doubles: a fake encoder (skips loading the real sentence-transformers
    model, which takes several seconds) and a throwaway Chroma directory
    (still real Chroma, since it's local and free - only the slow model
    load needs faking).
    """
    fake_embedding_service = EmbeddingService(encoder=fake_encoder)
    monkeypatch.setattr(documents_module, "embedding_service", fake_embedding_service)

    test_vector_store = VectorStoreService(persist_directory=tmp_path / "chroma_test")
    monkeypatch.setattr(documents_module, "vector_store_service", test_vector_store)

    return test_vector_store


@pytest.fixture()
def client() -> TestClient:
    return TestClient(app)


@pytest.fixture()
def valid_pdf_bytes() -> bytes:
    """A minimal, single-page, unencrypted PDF built at test time."""
    writer = PdfWriter()
    writer.add_blank_page(width=72, height=72)
    writer.add_metadata({"/Title": "PDF Test Title", "/Author": "PDF Test Author"})

    buffer = io.BytesIO()
    writer.write(buffer)
    return buffer.getvalue()


@pytest.fixture()
def encrypted_pdf_bytes() -> bytes:
    """A minimal PDF protected with a user password."""
    writer = PdfWriter()
    writer.add_blank_page(width=72, height=72)
    writer.encrypt(user_password="secret")

    buffer = io.BytesIO()
    writer.write(buffer)
    return buffer.getvalue()


@pytest.fixture()
def corrupted_pdf_bytes() -> bytes:
    """Bytes with a valid PDF signature but no parsable structure."""
    return b"%PDF-1.4\n" + b"this is not a real pdf body" * 20


@pytest.fixture()
def valid_docx_bytes() -> bytes:
    """A minimal DOCX with a couple of paragraphs, built at test time."""
    document = DocxDocument()
    document.core_properties.title = "DOCX Test Title"
    document.core_properties.author = "DOCX Test Author"
    document.add_paragraph("Knowledge base test document.")
    document.add_paragraph("Second paragraph.")

    buffer = io.BytesIO()
    document.save(buffer)
    return buffer.getvalue()


@pytest.fixture()
def corrupted_docx_bytes() -> bytes:
    """Bytes with a valid ZIP signature but no parsable DOCX structure."""
    return b"PK\x03\x04" + b"this is not a real docx body" * 20


@pytest.fixture()
def valid_pptx_bytes() -> bytes:
    """A minimal two-slide presentation, built at test time."""
    presentation = Presentation()
    presentation.core_properties.title = "PPTX Test Title"
    presentation.core_properties.author = "PPTX Test Author"
    layout = presentation.slide_layouts[1]

    slide_one = presentation.slides.add_slide(layout)
    slide_one.shapes.title.text = "Knowledge base test deck"
    slide_one.placeholders[1].text = "First slide content"

    slide_two = presentation.slides.add_slide(layout)
    slide_two.shapes.title.text = "Second slide"

    buffer = io.BytesIO()
    presentation.save(buffer)
    return buffer.getvalue()


@pytest.fixture()
def corrupted_pptx_bytes() -> bytes:
    """Bytes with a valid ZIP signature but no parsable PPTX structure."""
    return b"PK\x03\x04" + b"this is not a real pptx body" * 20
