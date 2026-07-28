import logging
from dataclasses import dataclass
from html.parser import HTMLParser
from pathlib import Path
from typing import Callable
from uuid import uuid4

from docx import Document as DocxDocument
from docx.opc.exceptions import PackageNotFoundError as DocxPackageNotFoundError
from fastapi import UploadFile
from pptx import Presentation
from pptx.exc import PackageNotFoundError as PptxPackageNotFoundError
from pypdf import PdfReader
from pypdf.errors import PdfReadError

from app.schemas.document import DocumentMetadata, DocumentRecord


logger = logging.getLogger(__name__)


@dataclass(frozen=True)
class DocumentFormat:
    """Metadata needed to validate an upload before it's saved or extracted."""

    extension: str
    allowed_content_types: frozenset[str]
    # Magic bytes the file must start with, or None if the format has none
    # reliable enough to check (e.g. plain text formats).
    signature: bytes | None = None


SUPPORTED_FORMATS: dict[str, DocumentFormat] = {
    ".pdf": DocumentFormat(
        extension=".pdf",
        allowed_content_types=frozenset({"application/pdf"}),
        signature=b"%PDF",
    ),
    ".txt": DocumentFormat(
        extension=".txt",
        allowed_content_types=frozenset({"text/plain"}),
    ),
    ".md": DocumentFormat(
        extension=".md",
        allowed_content_types=frozenset(
            {"text/markdown", "text/x-markdown", "text/plain"}
        ),
    ),
    ".docx": DocumentFormat(
        extension=".docx",
        allowed_content_types=frozenset(
            {
                "application/vnd.openxmlformats-officedocument"
                ".wordprocessingml.document"
            }
        ),
        # DOCX files are ZIP archives; every ZIP starts with this local
        # file header signature.
        signature=b"PK\x03\x04",
    ),
    ".pptx": DocumentFormat(
        extension=".pptx",
        allowed_content_types=frozenset(
            {
                "application/vnd.openxmlformats-officedocument"
                ".presentationml.presentation"
            }
        ),
        # PPTX is also a ZIP-based Office format.
        signature=b"PK\x03\x04",
    ),
    ".html": DocumentFormat(
        extension=".html",
        allowed_content_types=frozenset({"text/html"}),
        # No reliable magic bytes for HTML, same as the other text formats.
    ),
}


class _VisibleTextExtractor(HTMLParser):
    """Collects text nodes, skipping <script> and <style> content."""

    def __init__(self) -> None:
        super().__init__()
        self._chunks: list[str] = []
        self._skip_depth = 0

    def handle_starttag(self, tag: str, attrs: list[tuple[str, str | None]]) -> None:
        if tag in ("script", "style"):
            self._skip_depth += 1

    def handle_endtag(self, tag: str) -> None:
        if tag in ("script", "style") and self._skip_depth > 0:
            self._skip_depth -= 1

    def handle_data(self, data: str) -> None:
        if self._skip_depth == 0:
            self._chunks.append(data)

    def get_text(self) -> str:
        return "".join(self._chunks).strip()


class DocumentService:
    def __init__(self) -> None:
        self._extractors: dict[str, Callable[[Path], tuple[str, int]]] = {
            ".pdf": self._extract_pdf,
            ".txt": self._extract_plain_text,
            ".md": self._extract_plain_text,
            ".docx": self._extract_docx,
            ".pptx": self._extract_pptx,
            ".html": self._extract_html,
        }
        self._metadata_extractors: dict[str, Callable[[Path], DocumentMetadata]] = {
            ".pdf": self._extract_pdf_metadata,
            ".docx": self._extract_docx_metadata,
            ".pptx": self._extract_pptx_metadata,
            # TXT, Markdown, and HTML have no metadata standard to read
            # from, so they're deliberately absent here rather than
            # mapped to a no-op - extract_metadata treats a missing
            # entry as "no metadata available."
        }

    async def save_document(
        self,
        file: UploadFile,
        upload_directory: Path,
        extension: str,
    ) -> Path:
        """
        Save the uploaded file using a generated ID, preserving its extension.

        We do not use the original filename because:
        - multiple users may upload files with the same name,
        - filenames may contain unsafe characters,
        - generated IDs make documents easier to track.
        """
        document_id = str(uuid4())
        destination = upload_directory / f"{document_id}{extension}"

        content = await file.read()
        destination.write_bytes(content)

        logger.info(
            "File saved successfully | document_id=%s | path=%s",
            document_id,
            destination,
        )

        return destination

    def extract_text(self, file_path: Path) -> tuple[str, int]:
        """
        Extract text from a saved file, dispatching by its extension.

        Returns:
            A tuple containing:
            - the extracted text,
            - a page count (1 for formats with no real pagination concept).
        """
        extension = file_path.suffix.lower()
        extractor = self._extractors.get(extension)

        if extractor is None:
            raise ValueError(f"Unsupported file extension: {extension}")

        return extractor(file_path)

    def _extract_pdf(self, file_path: Path) -> tuple[str, int]:
        try:
            reader = PdfReader(file_path)

            # Encrypted PDFs require special handling.
            # For now, we reject them instead of attempting decryption.
            if reader.is_encrypted:
                raise ValueError("Encrypted PDF files are not supported.")

            page_texts: list[str] = []

            for page_number, page in enumerate(reader.pages, start=1):
                page_text = page.extract_text() or ""
                page_texts.append(page_text)

                logger.info(
                    "Page processed | page_number=%s | characters=%s",
                    page_number,
                    len(page_text),
                )

            extracted_text = "\n\n".join(page_texts)

            logger.info(
                "Text extraction completed | pages=%s | characters=%s",
                len(reader.pages),
                len(extracted_text),
            )

            return extracted_text, len(reader.pages)

        except PdfReadError as exc:
            logger.exception("PDF parsing failed | path=%s", file_path)
            raise ValueError("The uploaded file is not a readable PDF.") from exc

    def _extract_plain_text(self, file_path: Path) -> tuple[str, int]:
        try:
            text = file_path.read_text(encoding="utf-8")
        except UnicodeDecodeError as exc:
            logger.exception("Plain text decoding failed | path=%s", file_path)
            raise ValueError("The uploaded file is not valid UTF-8 text.") from exc

        logger.info(
            "Text extraction completed | characters=%s",
            len(text),
        )

        # Plain text formats have no real pagination concept.
        return text, 1

    def _extract_docx(self, file_path: Path) -> tuple[str, int]:
        try:
            document = DocxDocument(str(file_path))
        except DocxPackageNotFoundError as exc:
            logger.exception("DOCX parsing failed | path=%s", file_path)
            raise ValueError("The uploaded file is not a readable DOCX.") from exc

        # Paragraph text only for now - tables and embedded objects are
        # not extracted yet.
        text = "\n".join(paragraph.text for paragraph in document.paragraphs)

        logger.info(
            "Text extraction completed | characters=%s",
            len(text),
        )

        # DOCX has no reliably accessible page count without rendering it.
        return text, 1

    def _extract_pptx(self, file_path: Path) -> tuple[str, int]:
        try:
            presentation = Presentation(str(file_path))
        except PptxPackageNotFoundError as exc:
            logger.exception("PPTX parsing failed | path=%s", file_path)
            raise ValueError("The uploaded file is not a readable PPTX.") from exc

        slide_texts: list[str] = []

        for slide_number, slide in enumerate(presentation.slides, start=1):
            shape_texts = [
                shape.text_frame.text
                for shape in slide.shapes
                if shape.has_text_frame
            ]
            slide_text = "\n".join(text for text in shape_texts if text)
            slide_texts.append(slide_text)

            logger.info(
                "Slide processed | slide_number=%s | characters=%s",
                slide_number,
                len(slide_text),
            )

        extracted_text = "\n\n".join(slide_texts)

        logger.info(
            "Text extraction completed | slides=%s | characters=%s",
            len(presentation.slides),
            len(extracted_text),
        )

        # Unlike DOCX/TXT/MD, slides are a real pagination concept - use
        # the actual slide count instead of hardcoding 1.
        return extracted_text, len(presentation.slides)

    def _extract_html(self, file_path: Path) -> tuple[str, int]:
        try:
            html_content = file_path.read_text(encoding="utf-8")
        except UnicodeDecodeError as exc:
            logger.exception("HTML decoding failed | path=%s", file_path)
            raise ValueError("The uploaded file is not valid UTF-8 text.") from exc

        # HTMLParser is deliberately lenient about malformed markup (real
        # web pages are rarely perfectly valid HTML), so there is no
        # "corrupted HTML" failure mode the way there is for PDF/DOCX/PPTX -
        # any UTF-8 text will produce some extracted result.
        parser = _VisibleTextExtractor()
        parser.feed(html_content)
        text = parser.get_text()

        logger.info(
            "Text extraction completed | characters=%s",
            len(text),
        )

        # HTML has no real pagination concept.
        return text, 1

    def extract_metadata(self, file_path: Path) -> DocumentMetadata:
        """
        Extract format-intrinsic metadata (title, author, creation date)
        from a saved file, dispatching by its extension.

        This is best-effort: unlike extract_text, a failure here doesn't
        fail the whole upload. The file already parsed successfully once
        (extract_text succeeded before this is called), so a metadata
        extraction failure is treated as "no metadata available" rather
        than a reason to reject an otherwise-good document.
        """
        extension = file_path.suffix.lower()
        extractor = self._metadata_extractors.get(extension)

        if extractor is None:
            return DocumentMetadata()

        try:
            return extractor(file_path)
        except Exception:
            logger.exception(
                "Metadata extraction failed, continuing without it | path=%s",
                file_path,
            )
            return DocumentMetadata()

    def _extract_pdf_metadata(self, file_path: Path) -> DocumentMetadata:
        info = PdfReader(file_path).metadata

        if info is None:
            return DocumentMetadata()

        return DocumentMetadata(
            title=info.title or None,
            author=info.author or None,
            created_at=info.creation_date,
        )

    def _extract_docx_metadata(self, file_path: Path) -> DocumentMetadata:
        props = DocxDocument(str(file_path)).core_properties

        return DocumentMetadata(
            title=props.title or None,
            author=props.author or None,
            created_at=props.created,
        )

    def _extract_pptx_metadata(self, file_path: Path) -> DocumentMetadata:
        props = Presentation(str(file_path)).core_properties

        return DocumentMetadata(
            title=props.title or None,
            author=props.author or None,
            created_at=props.created,
        )

    def persist_metadata(
        self,
        record: DocumentRecord,
        upload_directory: Path,
    ) -> Path:
        """
        Persist a document's extracted text and metadata as a JSON sidecar
        file next to its saved file, named after the same document ID.
        """
        destination = upload_directory / f"{record.document_id}.json"
        destination.write_text(record.model_dump_json(indent=2))

        logger.info(
            "Metadata persisted | document_id=%s | path=%s",
            record.document_id,
            destination,
        )

        return destination
